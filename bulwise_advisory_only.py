import streamlit as st
import anthropic
import json
from datetime import datetime, timedelta
import hashlib
import sqlite3
import re
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import io

# Try to import streamlit-mermaid, fallback to code display if not available
try:
    from streamlit_mermaid import st_mermaid
    MERMAID_AVAILABLE = True
except ImportError:
    MERMAID_AVAILABLE = False

# Try to import reportlab for PDF generation
try:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Table, TableStyle
    from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_JUSTIFY
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

# Page configuration
st.set_page_config(
    page_title="Bulwise - Strategic AI Stack Advisory",
    page_icon="🎯",
    layout="wide",
    initial_sidebar_state="collapsed"
)

# ============================================================================
# RATE LIMITING - Cookie-Based
# ============================================================================

def get_browser_id():
    """Generate a unique browser ID for rate limiting."""
    if 'browser_id' not in st.session_state:
        import random
        st.session_state.browser_id = hashlib.md5(
            f"{datetime.now().timestamp()}{random.random()}".encode()
        ).hexdigest()
    return st.session_state.browser_id

def get_usage_data():
    """Get usage data from browser storage."""
    browser_id = get_browser_id()
    
    if 'usage_data' not in st.session_state:
        st.session_state.usage_data = {
            'browser_id': browser_id,
            'count': 0,
            'reset_date': datetime.now().isoformat(),
            'first_use': datetime.now().isoformat()
        }
    
    usage = st.session_state.usage_data
    
    # Check if we need to reset (monthly)
    reset_date = datetime.fromisoformat(usage['reset_date'])
    if datetime.now() > reset_date + timedelta(days=30):
        usage['count'] = 0
        usage['reset_date'] = datetime.now().isoformat()
        st.session_state.usage_data = usage
    
    return usage

def increment_usage():
    """Increment usage counter."""
    usage = get_usage_data()
    usage['count'] += 1
    st.session_state.usage_data = usage
    return usage['count']

def check_rate_limit():
    """Check if user has exceeded rate limit. Returns (allowed, count, reset_date)."""
    MAX_QUERIES_PER_MONTH = 3
    usage = get_usage_data()
    reset_date = datetime.fromisoformat(usage['reset_date'])
    next_reset = reset_date + timedelta(days=30)
    
    allowed = usage['count'] < MAX_QUERIES_PER_MONTH
    return allowed, usage['count'], next_reset

# ============================================================================
# ANALYTICS TRACKING
# ============================================================================

def init_analytics_db():
    """Initialize analytics database."""
    conn = sqlite3.connect('bulwise_analytics.db')
    c = conn.cursor()
    
    c.execute('''CREATE TABLE IF NOT EXISTS queries
                 (query_id TEXT PRIMARY KEY,
                  timestamp TEXT,
                  user_query TEXT,
                  query_length INTEGER,
                  detected_category TEXT,
                  budget_mentioned BOOLEAN,
                  team_size TEXT,
                  recommended_tools TEXT,
                  response_time REAL,
                  user_session_id TEXT,
                  rate_limit_hit BOOLEAN,
                  is_return_user BOOLEAN,
                  visit_count INTEGER,
                  days_since_last_visit REAL)''')
    
    c.execute('''CREATE TABLE IF NOT EXISTS user_sessions
                 (user_session_id TEXT PRIMARY KEY,
                  first_visit TEXT,
                  last_visit TEXT,
                  total_queries INTEGER,
                  total_visits INTEGER)''')
    
    c.execute('''CREATE TABLE IF NOT EXISTS feedback
                 (feedback_id TEXT PRIMARY KEY,
                  timestamp TEXT,
                  user_query TEXT,
                  sentiment TEXT,
                  feedback_text TEXT,
                  user_session_id TEXT)''')
    
    conn.commit()
    conn.close()

def detect_query_category(query):
    """Detect category from user query."""
    query_lower = query.lower()
    
    categories = {
        'video': ['video', 'youtube', 'tiktok', 'reels', 'editing'],
        'content': ['content', 'blog', 'article', 'writing', 'copy', 'social media'],
        'coding': ['code', 'coding', 'software', 'developer', 'programming', 'engineer'],
        'marketing': ['marketing', 'seo', 'email', 'campaign', 'ads'],
        'data': ['data', 'analytics', 'analysis', 'insights', 'reporting'],
        'design': ['design', 'graphic', 'ui', 'ux', 'creative'],
        'audio': ['audio', 'podcast', 'voice', 'music', 'sound'],
        'automation': ['automation', 'workflow', 'process', 'automate'],
        'customer_service': ['customer', 'support', 'service', 'chat']
    }
    
    for category, keywords in categories.items():
        if any(keyword in query_lower for keyword in keywords):
            return category
    
    return 'general'

def extract_budget(query):
    """Extract budget information from query."""
    budget_patterns = [r'\$\d+', r'\d+\s*dollars?', r'budget.*?\d+', r'cost', r'price', r'afford']
    query_lower = query.lower()
    
    for pattern in budget_patterns:
        if re.search(pattern, query_lower):
            return True
    return False

def extract_team_size(query):
    """Extract team size from query."""
    query_lower = query.lower()
    
    if any(word in query_lower for word in ['solo', 'freelance', 'individual', 'myself', 'i am']):
        return 'solo'
    elif any(word in query_lower for word in ['team of', 'people', 'members', 'employees']):
        numbers = re.findall(r'\d+', query)
        if numbers:
            size = int(numbers[0])
            if size <= 5:
                return 'small_team'
            elif size <= 20:
                return 'medium_team'
            else:
                return 'enterprise'
        return 'small_team'
    
    return 'unspecified'

def get_user_session_info():
    """Get or create user session info."""
    browser_id = get_browser_id()
    
    conn = sqlite3.connect('bulwise_analytics.db')
    c = conn.cursor()
    
    c.execute("SELECT * FROM user_sessions WHERE user_session_id = ?", (browser_id,))
    session = c.fetchone()
    
    if session:
        session_id, first_visit, last_visit, total_queries, total_visits = session
        
        last_visit_dt = datetime.fromisoformat(last_visit)
        days_since = (datetime.now() - last_visit_dt).total_seconds() / 86400
        
        is_return = True
        visit_count = total_visits + 1
        
        conn.close()
        return is_return, visit_count, days_since
    else:
        conn.close()
        return False, 1, 0.0

def update_user_session():
    """Update user session after query."""
    browser_id = get_browser_id()
    
    conn = sqlite3.connect('bulwise_analytics.db')
    c = conn.cursor()
    
    c.execute("SELECT * FROM user_sessions WHERE user_session_id = ?", (browser_id,))
    session = c.fetchone()
    
    if session:
        c.execute("""UPDATE user_sessions 
                     SET last_visit = ?, total_queries = total_queries + 1, total_visits = total_visits + 1
                     WHERE user_session_id = ?""",
                  (datetime.now().isoformat(), browser_id))
    else:
        c.execute("""INSERT INTO user_sessions VALUES (?, ?, ?, ?, ?)""",
                  (browser_id, datetime.now().isoformat(), datetime.now().isoformat(), 1, 1))
    
    conn.commit()
    conn.close()

def log_query(user_query, recommended_tools=None, response_time=None):
    """Log query to analytics database."""
    try:
        init_analytics_db()
        
        query_id = hashlib.md5(f"{user_query}{datetime.now().timestamp()}".encode()).hexdigest()
        timestamp = datetime.now().isoformat()
        query_length = len(user_query)
        category = detect_query_category(user_query)
        budget_mentioned = extract_budget(user_query)
        team_size = extract_team_size(user_query)
        browser_id = get_browser_id()
        
        allowed, count, _ = check_rate_limit()
        rate_limit_hit = not allowed
        
        is_return, visit_count, days_since = get_user_session_info()
        
        conn = sqlite3.connect('bulwise_analytics.db')
        c = conn.cursor()
        
        c.execute("""INSERT INTO queries VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)""",
                  (query_id, timestamp, user_query, query_length, category,
                   budget_mentioned, team_size, recommended_tools, response_time,
                   browser_id, rate_limit_hit, is_return, visit_count, days_since))
        
        conn.commit()
        conn.close()
        
        update_user_session()
        
    except Exception as e:
        st.error(f"Analytics logging error: {str(e)}")

def log_feedback(user_query, sentiment, feedback_text):
    """Log user feedback to analytics database."""
    try:
        feedback_id = hashlib.md5(f"{user_query}{datetime.now().timestamp()}".encode()).hexdigest()
        timestamp = datetime.now().isoformat()
        browser_id = get_browser_id()
        
        conn = sqlite3.connect('bulwise_analytics.db')
        c = conn.cursor()
        
        c.execute("""INSERT INTO feedback VALUES (?, ?, ?, ?, ?, ?)""",
                  (feedback_id, timestamp, user_query, sentiment, feedback_text, browser_id))
        
        conn.commit()
        conn.close()
        
    except Exception as e:
        st.error(f"Feedback logging error: {str(e)}")

# ============================================================================
# POWERPOINT GENERATION
# ============================================================================

def generate_professional_pptx(report_text, user_query):
    """Generate a professional PowerPoint presentation from the report with clean formatting."""
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme - Professional green/gray
    BRAND_GREEN = RGBColor(46, 125, 50)  # #2E7D32
    DARK_GRAY = RGBColor(51, 51, 51)
    LIGHT_GRAY = RGBColor(245, 245, 245)
    WHITE = RGBColor(255, 255, 255)
    
    def clean_text(text):
        """Remove markdown formatting for clean PowerPoint text."""
        # Remove bold markers
        text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
        # Remove italic markers
        text = re.sub(r'\*(.*?)\*', r'\1', text)
        # Remove headers
        text = re.sub(r'^#+\s+', '', text, flags=re.MULTILINE)
        # Remove bullet points
        text = re.sub(r'^\s*[-•]\s+', '', text, flags=re.MULTILINE)
        # Remove extra whitespace
        text = re.sub(r'\n\s*\n', '\n', text)
        return text.strip()
    
    def add_header_slide(slide, title_text, bg_color=BRAND_GREEN):
        """Add consistent header bar to slide."""
        header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
        header.fill.solid()
        header.fill.fore_color.rgb = bg_color
        header.line.color.rgb = bg_color
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
        tf = title_box.text_frame
        tf.text = title_text
        p = tf.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE
    
    # Parse sections from report
    sections = {}
    current_section = None
    current_content = []
    
    for line in report_text.split('\n'):
        if line.startswith('## '):
            if current_section:
                sections[current_section] = '\n'.join(current_content)
            current_section = line.replace('## ', '').strip()
            current_content = []
        else:
            current_content.append(line)
    
    if current_section:
        sections[current_section] = '\n'.join(current_content)
    
    # ============================================================================
    # SLIDE 1: TITLE SLIDE
    # ============================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_GRAY
    
    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
    tf = title_box.text_frame
    tf.text = "Strategic AI Implementation Advisory"
    p = tf.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = BRAND_GREEN
    p.alignment = PP_ALIGN.CENTER
    
    # User query subtitle
    query_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(2.5))
    tf = query_box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = 1  # Top align
    
    # Truncate if very long
    display_query = user_query if len(user_query) <= 300 else user_query[:300] + "..."
    tf.text = display_query
    
    p = tf.paragraphs[0]
    p.font.size = Pt(15)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.2
    
    # Date
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    tf = date_box.text_frame
    tf.text = f"Bulwise Advisory Report | {datetime.now().strftime('%B %d, %Y')}"
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # ============================================================================
    # SLIDE 2: EXECUTIVE SUMMARY
    # ============================================================================
    if 'EXECUTIVE SUMMARY' in sections:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_header_slide(slide, "Executive Summary")
        
        content = clean_text(sections['EXECUTIVE SUMMARY'])
        content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.2), Inches(8.5), Inches(5.8))
        tf = content_box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = 1  # Top align
        
        # Limit content to prevent overflow
        max_chars = 1200
        truncated_content = content[:max_chars] + ("..." if len(content) > max_chars else "")
        
        # Split into paragraphs
        for para in truncated_content.split('\n\n'):
            if para.strip():
                p = tf.add_paragraph() if len(tf.paragraphs) > 1 or tf.paragraphs[0].text else tf.paragraphs[0]
                p.text = para.strip()
                p.font.size = Pt(13)
                p.font.color.rgb = DARK_GRAY
                p.space_before = Pt(8)
                p.space_after = Pt(8)
                p.line_spacing = 1.2
    
    # ============================================================================
    # SLIDE 3: RECOMMENDED TOOLS - DETAILED
    # ============================================================================
    if 'STRATEGIC RECOMMENDATIONS' in sections:
        content = sections['STRATEGIC RECOMMENDATIONS']
        
        # Extract individual tool sections
        tool_sections = []
        current_tool = None
        current_tool_content = []
        
        for line in content.split('\n'):
            # Look for tool names (lines with ** at start and end on same line)
            if line.strip().startswith('**') and '**' in line[2:]:
                if current_tool:
                    tool_sections.append((current_tool, '\n'.join(current_tool_content)))
                # Extract tool name
                current_tool = line.strip().split('**')[1]
                current_tool_content = []
            elif current_tool:
                current_tool_content.append(line)
        
        if current_tool:
            tool_sections.append((current_tool, '\n'.join(current_tool_content)))
        
        # Create detailed slides for each tool
        for tool_name, tool_content in tool_sections[:6]:  # Max 6 tools
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # Truncate tool name if too long for title
            display_name = tool_name[:40] + "..." if len(tool_name) > 40 else tool_name
            add_header_slide(slide, f"Tool: {display_name}")
            
            # Clean and parse content
            clean_content = clean_text(tool_content)
            
            content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.2), Inches(8.5), Inches(5.8))
            tf = content_box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = 1  # Top align
            
            lines = [l for l in clean_content.split('\n') if l.strip()]
            char_count = 0
            max_chars = 1200  # Character limit per slide
            
            for line in lines:
                if char_count >= max_chars:
                    break
                
                # Truncate individual lines that are too long
                display_line = line.strip()
                if len(display_line) > 120:
                    display_line = display_line[:120] + "..."
                
                p = tf.add_paragraph() if len(tf.paragraphs) > 1 or tf.paragraphs[0].text else tf.paragraphs[0]
                p.text = display_line
                p.font.size = Pt(11)
                p.font.color.rgb = DARK_GRAY
                p.space_before = Pt(4)
                
                char_count += len(display_line)
    
    # ============================================================================
    # SLIDE: IMPLEMENTATION ROADMAP
    # ============================================================================
    if 'IMPLEMENTATION ROADMAP' in sections:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_header_slide(slide, "Implementation Roadmap")
        
        content = clean_text(sections['IMPLEMENTATION ROADMAP'])
        
        # Extract phases
        phases = []
        current_phase = None
        current_phase_content = []
        
        for line in content.split('\n'):
            if 'Phase' in line and ':' in line:
                if current_phase:
                    phases.append((current_phase, '\n'.join(current_phase_content)))
                current_phase = line.strip()
                current_phase_content = []
            elif current_phase:
                current_phase_content.append(line)
        
        if current_phase:
            phases.append((current_phase, '\n'.join(current_phase_content)))
        
        # Display phases in boxes
        y_start = Inches(1.5)
        max_phases = min(len(phases), 4)  # Max 4 phases to prevent overflow
        
        for idx, (phase_name, phase_content) in enumerate(phases[:max_phases]):
            # Phase box
            box = slide.shapes.add_shape(
                1,  # Rectangle
                Inches(0.75),
                y_start + (idx * Inches(1.3)),
                Inches(8.5),
                Inches(1.1)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = LIGHT_GRAY if idx % 2 else WHITE
            box.line.color.rgb = BRAND_GREEN
            box.line.width = Pt(2)
            
            # Phase title - truncate if too long
            tf = box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = 1  # Top
            p = tf.paragraphs[0]
            p.text = phase_name[:80]  # Truncate long titles
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = BRAND_GREEN
            p.space_after = Pt(6)
            
            # Phase content (first 2 lines only, truncated)
            content_lines = [l.strip() for l in phase_content.split('\n') if l.strip()][:2]
            for line in content_lines:
                p = tf.add_paragraph()
                # Truncate long lines
                display_line = line[:90] + "..." if len(line) > 90 else line
                p.text = display_line
                p.font.size = Pt(10)
                p.font.color.rgb = DARK_GRAY
                p.space_before = Pt(2)
    
    # ============================================================================
    # SLIDE: FINANCIAL ANALYSIS - WITH TABLE
    # ============================================================================
    if 'FINANCIAL ANALYSIS' in sections:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_header_slide(slide, "Financial Analysis")
        
        content = clean_text(sections['FINANCIAL ANALYSIS'])
        
        # Try to extract cost data for table
        monthly_costs = []
        total_monthly = None
        annual_total = None
        
        # Look for tool costs in format "Tool Name: $XX/month" or similar
        for line in content.split('\n'):
            # Extract monthly total
            if 'total monthly' in line.lower() and '$' in line:
                numbers = re.findall(r'\$(\d+(?:,\d{3})*(?:\.\d{2})?)', line)
                if numbers:
                    total_monthly = numbers[0].replace(',', '')
            
            # Extract annual total
            if 'annual' in line.lower() and '$' in line:
                numbers = re.findall(r'\$(\d+(?:,\d{3})*(?:\.\d{2})?)', line)
                if numbers:
                    annual_total = numbers[0].replace(',', '')
            
            # Extract individual tool costs (lines with $ and month/mo)
            if '$' in line and ('month' in line.lower() or '/mo' in line.lower()):
                # Try to extract tool name and cost
                parts = line.split(':')
                if len(parts) >= 2:
                    tool = parts[0].strip().strip('-•').strip()
                    cost_match = re.search(r'\$(\d+(?:,\d{3})*(?:\.\d{2})?)', parts[1])
                    if cost_match and len(monthly_costs) < 8:  # Max 8 tools
                        monthly_costs.append((tool, f"${cost_match.group(1)}"))
        
        # Create table if we have cost data
        if monthly_costs:
            rows = len(monthly_costs) + 2  # Tools + header + total
            cols = 2
            
            left = Inches(1)
            top = Inches(1.5)
            width = Inches(8)
            height = Inches(4.5)
            
            table = slide.shapes.add_table(rows, cols, left, top, width, height).table
            
            # Set column widths
            table.columns[0].width = Inches(5)
            table.columns[1].width = Inches(3)
            
            # Header row
            table.cell(0, 0).text = "Item"
            table.cell(0, 1).text = "Monthly Cost"
            
            for i in range(cols):
                cell = table.cell(0, i)
                cell.fill.solid()
                cell.fill.fore_color.rgb = BRAND_GREEN
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(14)
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
            
            # Tool rows
            for idx, (tool, cost) in enumerate(monthly_costs, 1):
                table.cell(idx, 0).text = tool[:60]  # Truncate if too long
                table.cell(idx, 1).text = cost
                
                for i in range(cols):
                    cell = table.cell(idx, i)
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(12)
                        paragraph.font.color.rgb = DARK_GRAY
                        if i == 1:  # Cost column
                            paragraph.alignment = PP_ALIGN.CENTER
            
            # Total row
            total_row = len(monthly_costs) + 1
            table.cell(total_row, 0).text = "TOTAL MONTHLY"
            table.cell(total_row, 1).text = f"${total_monthly}" if total_monthly else "See report"
            
            for i in range(cols):
                cell = table.cell(total_row, i)
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(13)
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = BRAND_GREEN
                    if i == 1:
                        paragraph.alignment = PP_ALIGN.CENTER
            
            # Add annual total below table if available
            if annual_total:
                text_box = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(8), Inches(0.5))
                tf = text_box.text_frame
                tf.text = f"Annual Projection: ${annual_total}"
                p = tf.paragraphs[0]
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = BRAND_GREEN
                p.alignment = PP_ALIGN.CENTER
        
        else:
            # Fallback: display as text if can't parse table
            content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.2), Inches(8.5), Inches(5.5))
            tf = content_box.text_frame
            tf.word_wrap = True
            
            lines = [l for l in content.split('\n') if l.strip()]
            for line in lines[:20]:  # Limit lines to prevent overflow
                p = tf.add_paragraph() if len(tf.paragraphs) > 1 or tf.paragraphs[0].text else tf.paragraphs[0]
                p.text = line.strip()[:100]  # Truncate long lines
                
                if ':' in line and len(line.split(':')[0]) < 40:
                    p.font.bold = True
                    p.font.size = Pt(12)
                    p.font.color.rgb = BRAND_GREEN
                else:
                    p.font.size = Pt(11)
                    p.font.color.rgb = DARK_GRAY
                
                p.space_before = Pt(4)
    
    # ============================================================================
    # SLIDE: RISK ASSESSMENT
    # ============================================================================
    risk_key = None
    for key in ['RISK ASSESSMENT & MITIGATION', 'RISK ASSESSMENT AND MITIGATION', 'RISK ASSESSMENT']:
        if key in sections:
            risk_key = key
            break
    
    if risk_key:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_header_slide(slide, "Risk Assessment")
        
        content = clean_text(sections[risk_key])
        
        content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.2), Inches(8.5), Inches(5.8))
        tf = content_box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = 1  # Top
        
        lines = [l for l in content.split('\n') if l.strip()]
        char_count = 0
        max_chars = 1000
        
        for line in lines:
            if char_count >= max_chars:
                break
            
            p = tf.add_paragraph() if len(tf.paragraphs) > 1 or tf.paragraphs[0].text else tf.paragraphs[0]
            # Truncate long lines
            display_line = line.strip()[:110] + ("..." if len(line.strip()) > 110 else "")
            p.text = display_line
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_GRAY
            p.space_before = Pt(4)
            
            char_count += len(display_line)
    
    # ============================================================================
    # SLIDE: ALTERNATIVE SCENARIOS
    # ============================================================================
    if 'ALTERNATIVE SCENARIOS' in sections:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_header_slide(slide, "Alternative Scenarios")
        
        content = clean_text(sections['ALTERNATIVE SCENARIOS'])
        
        content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.2), Inches(8.5), Inches(5.8))
        tf = content_box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = 1  # Top
        
        lines = [l for l in content.split('\n') if l.strip()]
        char_count = 0
        max_chars = 1000
        
        for line in lines:
            if char_count >= max_chars:
                break
            
            p = tf.add_paragraph() if len(tf.paragraphs) > 1 or tf.paragraphs[0].text else tf.paragraphs[0]
            # Truncate long lines
            display_line = line.strip()[:110] + ("..." if len(line.strip()) > 110 else "")
            p.text = display_line
            
            # Make scenario headers bold
            if 'Alternative' in display_line or 'Scenario' in display_line:
                p.font.bold = True
                p.font.size = Pt(12)
                p.font.color.rgb = BRAND_GREEN
            else:
                p.font.size = Pt(11)
                p.font.color.rgb = DARK_GRAY
            
            p.space_before = Pt(4)
            char_count += len(display_line)
    
    # ============================================================================
    # FINAL SLIDE: NEXT STEPS
    # ============================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BRAND_GREEN
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    tf.text = "Ready to Implement?"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(2))
    tf = contact_box.text_frame
    tf.text = "Get personalized implementation support\n\nhello@bulwise.io\nbulwise.io"
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = WHITE
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.space_before = Pt(10)
    
    # Save to BytesIO
    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    
    return pptx_io

# ============================================================================
# PDF GENERATION FUNCTION
# ============================================================================

def generate_professional_pdf(report_text, user_query):
    """Generate a professional PDF document from the report."""
    
    if not PDF_AVAILABLE:
        return None
    
    # Create PDF in memory
    pdf_buffer = io.BytesIO()
    
    # Create PDF document
    doc = SimpleDocTemplate(
        pdf_buffer,
        pagesize=letter,
        rightMargin=0.75*inch,
        leftMargin=0.75*inch,
        topMargin=1*inch,
        bottomMargin=0.75*inch
    )
    
    # Define colors
    BRAND_GREEN = colors.HexColor('#2E7D32')
    DARK_GRAY = colors.HexColor('#333333')
    LIGHT_GRAY = colors.HexColor('#F5F5F5')
    
    # Define styles
    styles = getSampleStyleSheet()
    
    # Custom styles
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=24,
        textColor=BRAND_GREEN,
        spaceAfter=30,
        alignment=TA_CENTER,
        fontName='Helvetica-Bold'
    )
    
    heading1_style = ParagraphStyle(
        'CustomHeading1',
        parent=styles['Heading1'],
        fontSize=16,
        textColor=BRAND_GREEN,
        spaceAfter=12,
        spaceBefore=20,
        fontName='Helvetica-Bold'
    )
    
    heading2_style = ParagraphStyle(
        'CustomHeading2',
        parent=styles['Heading2'],
        fontSize=14,
        textColor=BRAND_GREEN,
        spaceAfter=10,
        spaceBefore=15,
        fontName='Helvetica-Bold'
    )
    
    body_style = ParagraphStyle(
        'CustomBody',
        parent=styles['BodyText'],
        fontSize=11,
        textColor=DARK_GRAY,
        spaceAfter=12,
        alignment=TA_JUSTIFY,
        leading=14
    )
    
    bullet_style = ParagraphStyle(
        'CustomBullet',
        parent=styles['BodyText'],
        fontSize=10,
        textColor=DARK_GRAY,
        leftIndent=20,
        spaceAfter=6,
        leading=13
    )
    
    # Build PDF content
    content = []
    
    # Title page
    content.append(Spacer(1, 1*inch))
    content.append(Paragraph("Strategic AI Implementation Advisory", title_style))
    content.append(Spacer(1, 0.3*inch))
    
    # User query
    query_text = user_query if len(user_query) <= 500 else user_query[:500] + "..."
    content.append(Paragraph(query_text, body_style))
    content.append(Spacer(1, 0.3*inch))
    
    # Date
    date_text = f"Bulwise Advisory Report | {datetime.now().strftime('%B %d, %Y')}"
    content.append(Paragraph(date_text, body_style))
    content.append(PageBreak())
    
    # Parse and add report sections
    sections = {}
    current_section = None
    current_content = []
    
    for line in report_text.split('\n'):
        if line.startswith('## '):
            if current_section:
                sections[current_section] = '\n'.join(current_content)
            current_section = line.replace('## ', '').strip()
            current_content = []
        else:
            current_content.append(line)
    
    if current_section:
        sections[current_section] = '\n'.join(current_content)
    
    # Add each section to PDF
    for section_title, section_content in sections.items():
        # Skip mermaid code blocks
        if 'IMPLEMENTATION ROADMAP' in section_title:
            # Remove mermaid blocks for PDF
            section_content = re.sub(r'```\s*mermaid.*?```', '[Visual Timeline - See Online Report]', section_content, flags=re.DOTALL)
        
        # Section heading
        content.append(Paragraph(section_title, heading1_style))
        
        # Clean markdown from content
        clean_content = section_content
        # Remove bold markers but keep the text
        clean_content = re.sub(r'\*\*(.*?)\*\*', r'<b>\1</b>', clean_content)
        # Remove italic markers
        clean_content = re.sub(r'\*(.*?)\*', r'<i>\1</i>', clean_content)
        # Remove headers
        clean_content = re.sub(r'^###\s+(.+)$', r'<b>\1</b>', clean_content, flags=re.MULTILINE)
        
        # Split into paragraphs
        paragraphs = clean_content.split('\n\n')
        
        for para in paragraphs:
            if para.strip():
                # Handle bullet points
                if para.strip().startswith('-') or para.strip().startswith('•'):
                    lines = para.split('\n')
                    for line in lines:
                        if line.strip():
                            bullet_text = line.strip().lstrip('-•').strip()
                            content.append(Paragraph(f"• {bullet_text}", bullet_style))
                else:
                    # Regular paragraph
                    # Handle URLs
                    para_with_links = re.sub(
                        r'(https?://[^\s]+)',
                        r'<link href="\1" color="blue">\1</link>',
                        para.strip()
                    )
                    content.append(Paragraph(para_with_links, body_style))
        
        content.append(Spacer(1, 0.2*inch))
    
    # Build PDF
    doc.build(content)
    
    # Get PDF data
    pdf_buffer.seek(0)
    return pdf_buffer

def validate_report_quality(report_text):
    """
    Validate report quality and return quality score with specific feedback.
    Returns: (score, issues_list, warnings_list)
    """
    
    issues = []
    warnings = []
    score = 100
    
    # Required sections
    required_sections = [
        'EXECUTIVE SUMMARY',
        'STRATEGIC RECOMMENDATIONS',
        'IMPLEMENTATION ROADMAP',
        'FINANCIAL ANALYSIS',
        'RISK ASSESSMENT'
    ]
    
    missing_sections = []
    for section in required_sections:
        if section not in report_text:
            missing_sections.append(section)
            score -= 15
    
    if missing_sections:
        issues.append(f"Missing critical sections: {', '.join(missing_sections)}")
    
    # Check for tool recommendations with URLs
    tool_count = len(re.findall(r'\*\*([^*]+)\*\*\s*\([^)]+\)', report_text))
    url_count = len(re.findall(r'https?://[^\s\)]+', report_text))
    
    if tool_count == 0:
        issues.append("No tools recommended")
        score -= 20
    elif tool_count < 3:
        warnings.append(f"Only {tool_count} tools recommended (recommend 3-5)")
        score -= 5
    
    if tool_count > 0 and url_count < tool_count:
        warnings.append(f"{tool_count - url_count} tools missing URLs")
        score -= 10
    
    # Check section lengths
    sections = {}
    current_section = None
    current_content = []
    
    for line in report_text.split('\n'):
        if line.startswith('## '):
            if current_section:
                sections[current_section] = '\n'.join(current_content)
            current_section = line.replace('## ', '').strip()
            current_content = []
        else:
            current_content.append(line)
    
    if current_section:
        sections[current_section] = '\n'.join(current_content)
    
    # Check Executive Summary length
    if 'EXECUTIVE SUMMARY' in sections:
        exec_summary_len = len(sections['EXECUTIVE SUMMARY'].strip())
        if exec_summary_len < 200:
            warnings.append(f"Executive Summary too brief ({exec_summary_len} chars, recommend 300+)")
            score -= 5
    
    # Check for cost information
    if 'FINANCIAL ANALYSIS' in sections:
        financial_text = sections['FINANCIAL ANALYSIS']
        if '$' not in financial_text:
            issues.append("Financial Analysis missing specific costs")
            score -= 10
        
        # Check for monthly/annual costs
        if 'month' not in financial_text.lower():
            warnings.append("Financial Analysis missing monthly breakdown")
            score -= 5
    
    # Check for implementation timeline
    if 'IMPLEMENTATION ROADMAP' in sections:
        roadmap = sections['IMPLEMENTATION ROADMAP']
        phase_count = len(re.findall(r'Phase \d+', roadmap))
        if phase_count == 0:
            warnings.append("Implementation Roadmap missing phase structure")
            score -= 5
        elif phase_count < 3:
            warnings.append(f"Only {phase_count} phases (recommend 3-4)")
    
    # Check for risk mitigation
    if 'RISK ASSESSMENT' in sections:
        risk_text = sections['RISK ASSESSMENT']
        if 'mitigation' not in risk_text.lower() and 'mitigate' not in risk_text.lower():
            warnings.append("Risk Assessment missing mitigation strategies")
            score -= 5
    
    # Ensure score doesn't go below 0
    score = max(0, score)
    
    return score, issues, warnings

def display_quality_badge(score):
    """Display quality score as colored badge."""
    
    if score >= 90:
        color = "#2E7D32"  # Green
        label = "Excellent"
        emoji = "🏆"
    elif score >= 75:
        color = "#558B2F"  # Light green
        label = "Good"
        emoji = "✅"
    elif score >= 60:
        color = "#FFA000"  # Orange
        label = "Acceptable"
        emoji = "⚠️"
    else:
        color = "#D32F2F"  # Red
        label = "Needs Improvement"
        emoji = "❌"
    
    return f"""
    <div style="background: {color}; color: white; padding: 0.5rem 1rem; border-radius: 8px; 
                display: inline-block; font-weight: bold; margin: 1rem 0;">
        {emoji} Quality Score: {score}/100 - {label}
    </div>
    """

def render_report_with_mermaid(report_text):
    """Render report text with Mermaid diagrams and tool logos."""
    
    # Split by mermaid blocks
    pattern = r'```\s*mermaid\s*\n(.*?)```'
    parts = re.split(pattern, report_text, flags=re.DOTALL | re.IGNORECASE)
    
    for i, part in enumerate(parts):
        if i % 2 == 0:
            # Regular markdown content - enhance with logos
            if part.strip():
                # Add logos to tool recommendations
                enhanced_part = add_tool_logos(part)
                st.markdown(enhanced_part, unsafe_allow_html=True)
        else:
            # Mermaid diagram content
            mermaid_code = part.strip()
            
            if MERMAID_AVAILABLE:
                # Use streamlit-mermaid library
                try:
                    st_mermaid(mermaid_code, height=500)
                except Exception as e:
                    # Fallback to styled code block
                    st.markdown("### 📊 Visual Timeline")
                    st.code(mermaid_code, language="mermaid")
            else:
                # Library not available - show as styled code
                st.markdown("### 📊 Implementation Timeline")
                st.info("💡 Interactive diagram - see details below")
                
                # Parse and display in a more readable format
                lines = mermaid_code.split('\n')
                for line in lines:
                    if line.strip() and not line.strip().startswith('graph'):
                        # Clean up the mermaid syntax
                        cleaned = line.replace('-->', '→').replace('[', '**').replace(']', '**')
                        st.markdown(f"- {cleaned.strip()}")

def add_tool_logos(text):
    """Add logos to tool names in markdown text."""
    
    # Pattern to find tool recommendations with websites
    # Looks for: **ToolName** followed by website URL
    pattern = r'\*\*([^*]+)\*\*\s*\([^)]+\)([^#]*?)(?:Official Website|Website):\s*\*?\*?(https?://[^\s\)]+)'
    
    def replace_with_logo(match):
        tool_name = match.group(1).strip()
        content = match.group(2)
        url = match.group(3).strip().rstrip('*').rstrip(')')
        
        # Extract domain for favicon
        domain_match = re.search(r'https?://(?:www\.)?([^/]+)', url)
        if domain_match:
            domain = domain_match.group(1)
            # Use DuckDuckGo's favicon service (free, no API key needed)
            logo_url = f"https://icons.duckduckgo.com/ip3/{domain}.ico"
            
            # Return enhanced HTML with logo
            return f'''
<div style="border-left: 4px solid #2E7D32; padding-left: 1rem; margin: 1.5rem 0;">
    <div style="display: flex; align-items: center; gap: 0.75rem; margin-bottom: 0.5rem;">
        <img src="{logo_url}" width="24" height="24" style="border-radius: 4px;" onerror="this.style.display='none'">
        <strong style="font-size: 1.1rem; color: #2E7D32;">{tool_name}</strong>
    </div>
    {content}
    <p><strong>Official Website:</strong> <a href="{url}" target="_blank" style="color: #2E7D32;">{url}</a></p>
</div>
'''
        return match.group(0)
    
    enhanced = re.sub(pattern, replace_with_logo, text, flags=re.DOTALL)
    return enhanced

# ============================================================================
# LOAD DATABASE
# ============================================================================

@st.cache_data
def load_database():
    """Load AI tools database."""
    try:
        with open('ai_tools_complete.json', 'r') as f:
            return json.load(f)
    except FileNotFoundError:
        st.error("Database file not found. Please ensure ai_tools_complete.json is in the app directory.")
        return {'ai_tools': []}

# ============================================================================
# CUSTOM CSS
# ============================================================================

st.markdown("""
    <style>
    /* Clean, minimal styling */
    .main {
        padding-top: 2rem;
    }
    
    /* Progress bar styling */
    .stProgress > div > div > div > div {
        background-color: #4CAF50;
    }
    
    /* Button styling */
    .stButton > button {
        background-color: #2E7D32;
        color: white;
        border-radius: 8px;
        padding: 0.75rem 2rem;
        font-weight: 600;
        border: none;
        transition: all 0.3s ease;
    }
    
    .stButton > button:hover {
        background-color: #1B5E20;
        box-shadow: 0 4px 12px rgba(0,0,0,0.15);
    }
    
    /* Text input styling */
    .stTextArea textarea {
        border-radius: 8px;
        border: 2px solid #e0e0e0;
        padding: 1rem;
    }
    
    /* Report styling */
    .whitepaper-section {
        background: white;
        padding: 2rem;
        border-radius: 12px;
        box-shadow: 0 2px 8px rgba(0,0,0,0.1);
        margin-top: 2rem;
        /* Enable copy/paste */
        user-select: text !important;
        -webkit-user-select: text !important;
        -moz-user-select: text !important;
        -ms-user-select: text !important;
    }
    
    .whitepaper-section * {
        user-select: text !important;
    }
    
    .whitepaper-header {
        font-size: 1.8rem;
        font-weight: 700;
        color: #1a1a1a;
        margin-bottom: 1.5rem;
        padding-bottom: 1rem;
        border-bottom: 3px solid #2E7D32;
    }
    
    /* Hide Streamlit branding */
    #MainMenu {visibility: hidden;}
    footer {visibility: hidden;}
    
    /* Responsive */
    @media (max-width: 768px) {
        .main {
            padding: 1rem;
        }
    }
    </style>
""", unsafe_allow_html=True)

# ============================================================================
# MAIN APP - ADVISORY ONLY
# ============================================================================

# Initialize analytics
init_analytics_db()

# Load database
database = load_database()

# Get API key from secrets or user input
try:
    api_key = st.secrets.get("ANTHROPIC_API_KEY", None)
except Exception:
    api_key = None

if not api_key:
    st.warning("⚠️ API key required")
    api_key = st.text_input(
        "Enter your Anthropic API Key:",
        type="password",
        help="Get your API key at console.anthropic.com"
    )
    
    if not api_key:
        st.info("💡 **No API key?** Get one at [console.anthropic.com](https://console.anthropic.com)")
        st.stop()

# ============================================================================
# HIDDEN ADMIN SECTION (Password Protected)
# ============================================================================

# Add small password input in sidebar (only visible if sidebar expanded)
with st.sidebar:
    admin_password = st.text_input("🔒", type="password", key="admin_pass", help="Admin access")
    
    if admin_password == "bulwise2024":  # Change this password!
        st.success("✅ Admin access granted")
        
        # Show admin panel
        st.markdown("---")
        st.markdown("### 📊 Admin Panel")
        
        # Download database button
        if st.button("💾 Download Analytics Database"):
            try:
                with open('bulwise_analytics.db', 'rb') as f:
                    st.download_button(
                        label="⬇️ Click to Download",
                        data=f,
                        file_name='bulwise_analytics.db',
                        mime='application/octet-stream'
                    )
            except FileNotFoundError:
                st.warning("No database file found yet")
        
        # Show quick stats
        try:
            conn = sqlite3.connect('bulwise_analytics.db')
            c = conn.cursor()
            
            # Total queries
            c.execute("SELECT COUNT(*) FROM queries")
            total = c.fetchone()[0]
            st.metric("Total Queries", total)
            
            # Recent queries
            c.execute("SELECT user_query, timestamp FROM queries ORDER BY timestamp DESC LIMIT 5")
            recent = c.fetchall()
            
            if recent:
                st.markdown("**Recent Queries:**")
                for query, timestamp in recent:
                    st.text(f"{timestamp[:10]}: {query[:50]}...")
            
            conn.close()
        except Exception as e:
            st.info("No analytics data yet")

# Header
st.markdown("# 🎯 Strategic AI Stack Advisory")
st.markdown("### Get professional recommendations tailored to your needs")
st.markdown("---")

st.markdown("<br>", unsafe_allow_html=True)

# Initialize step tracking
if 'workflow_step' not in st.session_state:
    st.session_state.workflow_step = 1
if 'initial_query' not in st.session_state:
    st.session_state.initial_query = ''
if 'clarifying_answers' not in st.session_state:
    st.session_state.clarifying_answers = {}

# STEP 1: Initial Query
if st.session_state.workflow_step == 1:
    # ChatGPT Differentiator
    st.markdown("""
    <div style='background: #f0f7ff; padding: 1.2rem; border-radius: 8px; border-left: 4px solid #2E7D32; margin-bottom: 1.5rem;'>
        <strong>🎯 Why Bulwise?</strong> Turn your AI question into a boardroom-ready strategic advisory report with implementation details, not just suggestions.
    </div>
    """, unsafe_allow_html=True)
    
    st.markdown("""
    <div style='background: #f8f9fa; padding: 1.5rem; border-radius: 8px; border-left: 4px solid #2E7D32; margin-bottom: 1.5rem;'>
        <h4 style='margin-top: 0; color: #2E7D32;'>💡 For Best Results, Include:</h4>
        <ul style='margin-bottom: 0;'>
            <li>What you're trying to accomplish</li>
            <li>Team size and budget (if relevant)</li>
            <li>Current blockers or constraints</li>
            <li>Timeline or urgency</li>
        </ul>
    </div>
    """, unsafe_allow_html=True)
    
    # Quick start examples - BEFORE TEXT AREA
    st.markdown("**Quick Start Examples:**")
    col1, col2, col3, col4 = st.columns(4)
    
    with col1:
        if st.button("📹 Video\nProduction", use_container_width=True, key="ex_video"):
            st.session_state.initial_query = "I run a digital agency and need to produce 10-15 professional video content pieces monthly for enterprise clients. Need AI tools for script writing, video generation, voiceovers, and editing. Budget: $200-300/month. Team of 3 people with mixed technical skills."
            st.rerun()
    
    with col2:
        if st.button("✍️ Content\nMarketing", use_container_width=True, key="ex_content"):
            st.session_state.initial_query = "E-commerce business needing to scale content creation: product descriptions, blog posts, email campaigns, and social media. Team of 2 marketers. Budget: $75-100/month."
            st.rerun()
    
    with col3:
        if st.button("💻 Software\nDevelopment", use_container_width=True, key="ex_dev"):
            st.session_state.initial_query = "Software development team of 5 engineers working on full-stack web applications. Need AI assistants for code generation, debugging, documentation, and code review. Budget: $100-150/month total."
            st.rerun()
    
    with col4:
        if st.button("🎙️ Podcast\nProduction", use_container_width=True, key="ex_podcast"):
            st.session_state.initial_query = "Starting a weekly B2B podcast. Need tools for recording, editing, transcription, show notes generation, and distribution. Solo operation with limited technical background. Budget: $50-75/month."
            st.rerun()
    
    st.markdown("<br>", unsafe_allow_html=True)
    
    # Text area - Force display of current session state by NOT using key parameter
    user_query = st.text_area(
        "**📝 Describe Your Situation:**",
        value=st.session_state.initial_query,
        height=180,
        placeholder="Example: I'm a marketing consultant working with 5-10 small business clients simultaneously. I need to create weekly content (blog posts, social media, email campaigns) efficiently. Budget is $100-150/month. I have intermediate technical skills and need tools that integrate well together. Timeline: implement within 2 weeks.",
        help="Click 'Continue →' button below when ready to proceed"
    )
    
    # Always update from user input (allows manual editing)
    if user_query != st.session_state.initial_query:
        st.session_state.initial_query = user_query
    
    st.markdown("---")
    
    # Continue button for Step 1
    col1, col2 = st.columns([3, 1])
    with col2:
        continue_button = st.button(
            "Continue →",
            type="primary",
            use_container_width=True,
            disabled=not st.session_state.initial_query.strip()
        )
        
        if continue_button:
            st.session_state.workflow_step = 2
            st.rerun()

# STEP 2: Clarifying Questions
elif st.session_state.workflow_step == 2:
    st.markdown("### 🔍 Refine Your Requirements")
    st.markdown("*These optional details help us provide more precise recommendations*")
    st.markdown("<br>", unsafe_allow_html=True)
    
    # Show what they entered
    with st.expander("📝 Your Initial Request", expanded=False):
        st.write(st.session_state.initial_query)
    
    st.markdown("---")
    
    # Group 1: Team & Scale
    st.markdown("#### 👥 Team & Scale")
    col1, col2 = st.columns(2)
    
    with col1:
        team_size = st.selectbox(
            "**Team Size:**",
            ["Not specified", "Individual (just me)", "Small team (2-10)", "Medium team (11-50)", "Large team (50+)"],
            index=0,
            key="team_size_select",
            help="Helps us recommend tools with appropriate collaboration features and pricing tiers"
        )
    
    with col2:
        use_case_type = st.selectbox(
            "**Primary Use Case:**",
            ["Not specified", "Content creation", "Software development", "Marketing & sales", "Data analysis", "Customer support", "Design & creative", "Other"],
            index=0,
            key="use_case_select",
            help="Helps us prioritize tools relevant to your specific workflow"
        )
    
    st.markdown("---")
    
    # Group 2: Budget & Timeline
    st.markdown("#### 💰 Budget & Timeline")
    col1, col2 = st.columns(2)
    
    with col1:
        budget = st.selectbox(
            "**Monthly Budget Range:**",
            ["Not specified", "Under $50/month", "$50-$200/month", "$200-$500/month", "$500-$1,000/month", "$1,000+/month"],
            index=0,
            key="budget_select",
            help="Your budget range helps us recommend cost-effective solutions that fit your constraints"
        )
    
    with col2:
        timeline = st.selectbox(
            "**Implementation Timeline:**",
            ["Not specified", "Urgent (need now)", "Within 2 weeks", "Within 1 month", "Within 3 months", "Flexible timeline"],
            index=0,
            key="timeline_select",
            help="Urgency affects complexity of tools we recommend (simpler for urgent needs)"
        )
    
    st.markdown("---")
    
    # Group 3: Technical Preferences
    st.markdown("#### ⚙️ Technical Preferences")
    col1, col2 = st.columns(2)
    
    with col1:
        experience = st.selectbox(
            "**Technical Experience:**",
            ["Not specified", "Beginner (need easy tools)", "Intermediate (comfortable with tech)", "Advanced (can handle complexity)"],
            index=0,
            key="experience_select",
            help="Determines the learning curve we factor into recommendations"
        )
    
    with col2:
        integration_need = st.selectbox(
            "**Integration Requirements:**",
            ["Not specified", "Must integrate with existing tools", "Standalone tools preferred", "Either works"],
            index=0,
            key="integration_select",
            help="Helps us prioritize tools with strong API support and integrations if needed"
        )
    
    st.markdown("---")
    
    # Group 4: Specific Constraints (Optional)
    with st.expander("🎯 Additional Constraints (Optional)", expanded=False):
        st.markdown("*These help us tailor recommendations even further*")
        
        data_privacy = st.selectbox(
            "**Data Privacy Requirements:**",
            ["Not specified", "Standard cloud services OK", "Need data residency options", "Must be on-premise/self-hosted"],
            index=0,
            key="privacy_select",
            help="Important for regulated industries or sensitive data"
        )
        
        existing_tools = st.text_input(
            "**Current Tools (if any):**",
            placeholder="e.g., Slack, Google Workspace, Salesforce",
            key="existing_tools_input",
            help="Helps us suggest tools that integrate well with your current stack"
        )
        
        must_have_features = st.text_area(
            "**Must-Have Features:**",
            placeholder="e.g., Mobile app, offline mode, custom branding",
            height=80,
            key="must_have_input",
            help="Specific features that are non-negotiable for your use case"
        )
    
    # Save all answers
    st.session_state.clarifying_answers = {
        'team_size': team_size if team_size != "Not specified" else None,
        'use_case_type': use_case_type if use_case_type != "Not specified" else None,
        'budget': budget if budget != "Not specified" else None,
        'timeline': timeline if timeline != "Not specified" else None,
        'experience': experience if experience != "Not specified" else None,
        'integration_need': integration_need if integration_need != "Not specified" else None,
        'data_privacy': data_privacy if data_privacy != "Not specified" else None,
        'existing_tools': existing_tools if existing_tools.strip() else None,
        'must_have_features': must_have_features if must_have_features.strip() else None
    }
    
    st.markdown("<br>", unsafe_allow_html=True)
    
    # Show what's been specified
    specified_count = sum(1 for v in st.session_state.clarifying_answers.values() if v)
    if specified_count > 0:
        st.info(f"✅ You've specified {specified_count} details to refine your recommendations")
    
    st.markdown("<br>", unsafe_allow_html=True)
    
    # Navigation buttons
    col1, col2, col3 = st.columns([1, 1, 1])
    
    with col1:
        if st.button("← Back", use_container_width=True):
            st.session_state.workflow_step = 1
            st.rerun()
    
    with col2:
        if st.button("Skip All →", use_container_width=True):
            st.session_state.clarifying_answers = {}
            st.session_state.workflow_step = 3
            st.rerun()
    
    with col3:
        if st.button("Continue →", type="primary", use_container_width=True):
            st.session_state.workflow_step = 3
            st.rerun()

# STEP 3: Generate Report
elif st.session_state.workflow_step == 3:
    # Build complete query with clarifying answers
    user_query = st.session_state.initial_query
    
    if any(st.session_state.clarifying_answers.values()):
        clarifications = []
        
        # Core details
        if st.session_state.clarifying_answers.get('team_size'):
            clarifications.append(f"Team size: {st.session_state.clarifying_answers['team_size']}")
        if st.session_state.clarifying_answers.get('use_case_type'):
            clarifications.append(f"Primary use case: {st.session_state.clarifying_answers['use_case_type']}")
        if st.session_state.clarifying_answers.get('budget'):
            clarifications.append(f"Budget: {st.session_state.clarifying_answers['budget']}")
        if st.session_state.clarifying_answers.get('timeline'):
            clarifications.append(f"Timeline: {st.session_state.clarifying_answers['timeline']}")
        if st.session_state.clarifying_answers.get('experience'):
            clarifications.append(f"Technical experience: {st.session_state.clarifying_answers['experience']}")
        
        # Technical preferences
        if st.session_state.clarifying_answers.get('integration_need'):
            clarifications.append(f"Integration preference: {st.session_state.clarifying_answers['integration_need']}")
        if st.session_state.clarifying_answers.get('data_privacy'):
            clarifications.append(f"Data privacy: {st.session_state.clarifying_answers['data_privacy']}")
        
        # Specific details
        if st.session_state.clarifying_answers.get('existing_tools'):
            clarifications.append(f"Current tools: {st.session_state.clarifying_answers['existing_tools']}")
        if st.session_state.clarifying_answers.get('must_have_features'):
            clarifications.append(f"Must-have features: {st.session_state.clarifying_answers['must_have_features']}")
        
        user_query += "\n\nAdditional context: " + "; ".join(clarifications)
    
    # Show summary
    st.markdown("### 📋 Review Your Request")
    with st.expander("View full request", expanded=True):
        st.write(user_query)
    
    # Back button
    if st.button("← Edit Request"):
        st.session_state.workflow_step = 1
        st.rerun()
    
    st.markdown("---")
    
    # Check rate limit and display status
    allowed, query_count, next_reset = check_rate_limit()
    
    # Display usage status
    col1, col2 = st.columns([3, 1])
    with col1:
        if allowed:
            remaining = 3 - query_count
            if remaining == 3:
                st.info(f"ℹ️ Free tier: {remaining} advisory reports remaining this month")
            elif remaining == 2:
                st.info(f"ℹ️ {remaining} advisory reports remaining this month")
            else:
                st.warning(f"⚠️ Only {remaining} advisory report remaining this month. Resets on {next_reset.strftime('%B %d, %Y')}")
        else:
            st.error(f"🚫 Monthly limit reached (3 free reports). Resets on {next_reset.strftime('%B %d, %Y')}")
            st.info("💡 Want unlimited reports? Contact: hello@bulwise.io")
    
    with col2:
        generate_button = st.button(
            "📊 Generate Advisory",
            type="primary",
            use_container_width=True,
            disabled=not allowed or not user_query.strip()
        )
    
    # Generate advisory report
    if generate_button and user_query.strip():
        
        with st.spinner(""):
            try:
                client = anthropic.Anthropic(api_key=api_key)
                
                # Create progress tracking
                progress_container = st.container()
                with progress_container:
                    progress_bar = st.progress(0)
                    status_text = st.empty()
                
                status_text.text("🔍 Analyzing your requirements... (10%)")
                progress_bar.progress(10)
                
                # System prompt - UPDATED with user validation, data sources upfront, no emojis, visualizations, enhanced costs
                system_prompt = f"""You are Bulwise, a strategic AI implementation advisory service. Generate professional, data-driven reports for clients seeking AI tool recommendations.

You have access to:
1. A comprehensive database of {len(database['ai_tools'])} AI tools with detailed specifications
2. Web search capability to find current market research and trends

CRITICAL: Format your response as a professional whitepaper-style advisory report with the following structure (IN THIS EXACT ORDER):

---

## ANALYST NOTE
[Brief professional note (2-3 sentences) welcoming the client, emphasizing the data-driven strategic nature of this recommendation, and expressing willingness to answer follow-up questions. Keep warm and professional. NO EMOJIS.]

---

## REQUIREMENTS VALIDATION
[NEW - #4: Before proceeding with analysis, confirm your understanding of the client's needs]

Based on your request, I understand you are seeking to:
- **Primary Goal:** [State what you understand they want to accomplish]
- **Team Context:** [Solo, small team, enterprise - based on their input]
- **Budget Range:** [If mentioned, confirm. If not mentioned, note "Budget not specified - recommendations will span multiple price points"]
- **Timeline:** [If mentioned, confirm. If not, note "Timeline not specified"]
- **Key Constraints:** [Any blockers or limitations they mentioned]

If any of these assumptions are incorrect, please let me know. Proceeding with analysis based on these parameters.

---

## EXECUTIVE SUMMARY
[2-3 sentences providing high-level overview of the recommendation and expected outcomes]

**Data Foundation:** [NEW - #5: Move data sources upfront]
This analysis draws from {len(database['ai_tools'])} AI tools across 15+ categories in our proprietary database, supplemented by current market research from Gartner, Forrester, IDC, and industry-specific sources. All statistics and benchmarks are cited with sources.

---

## METHODOLOGY & ANALYSIS

**Requirements Analysis:**
[Summarize the client's stated needs, constraints, and objectives]

**Evaluation Criteria:**
[List the key factors considered in tool selection: budget alignment, technical requirements, workflow integration, scalability, etc.]

---

## STRATEGIC RECOMMENDATIONS

### Recommended Tools
[FIRST: List the 3-5 recommended tool NAMES ONLY as a simple bulleted list for quick reference]

Example format:
- **ChatGPT** (Content Generation)
- **Midjourney** (Visual Design)
- **Claude** (Strategic Analysis)

### Detailed Technology Stack Analysis
[NOW: Provide the full detailed analysis for each tool]

For each tool, provide:
- **Tool Name** (Category)
- **Strategic Rationale:** Why this tool fits the client's needs (2-3 sentences with specific data points, referencing market position if available)
- **Key Capabilities:** Bullet list of relevant features
- **Pricing:** Specific tier recommendation with justification
- **Official Website:** [REQUIRED - Always include the tool's official website URL from the database. Format as: https://toolname.com]

**CRITICAL:** Every tool recommendation MUST include its official website URL. Look up the URL in the database and include it prominently.

### Technology Stack Visualization
[NEW - #14: ASCII diagram showing how tools connect]

```
Data Flow Architecture:

[Input Source] → [Primary Tool] → [Enhancement Tool] → [Output/Storage]
                      ↓
                [Supporting Tool]
                      ↓
                [Quality Control]

Example:
User Content → ChatGPT (Generation) → Grammarly (Polish) → Notion (Storage)
                      ↓
                  Midjourney (Visuals)
                      ↓
                  Canva (Assembly)
```

### Tool Integration Workflow
[Provide step-by-step explanation of how these tools work together in practice, including data flow and handoffs between systems]

---

## IMPLEMENTATION ROADMAP

### Visual Timeline
[NEW - #15: Mermaid diagram]

```mermaid
graph LR
    A[Week 1-2: Foundation] --> B[Week 3-4: Integration]
    B --> C[Month 2: Optimization]
    C --> D[Month 3+: Scaling]
    
    A --> A1[Tool Setup]
    A --> A2[Team Training]
    B --> B1[Workflow Integration]
    B --> B2[Testing]
    C --> C1[Process Refinement]
    C --> C2[Advanced Features]
```

### Phase 1: Foundation (Week 1-2)
[Initial setup steps and quick wins]

### Phase 2: Integration (Week 3-4)
[Connecting tools and establishing workflows]

### Phase 3: Optimization (Month 2+)
[Refinement and advanced feature adoption]

---

## FINANCIAL ANALYSIS

### Detailed Cost Breakdown
[NEW - #16: Enhanced with per-tool details]

**Monthly Costs by Tool:**

| Tool | Recommended Tier | Monthly Cost | Annual Cost | Notes |
|------|-----------------|--------------|-------------|-------|
| [Tool 1] | [Tier name] | $XX | $XXX | [Why this tier] |
| [Tool 2] | [Tier name] | $XX | $XXX | [Why this tier] |
| [Tool 3] | [Tier name] | $XX | $XXX | [Why this tier] |

**Cost Summary:**
- **Total Monthly Investment:** $[X]
- **Annual Projection:** $[Y]
- **Cost Per User (if team):** $[Z]

**Payment Options:**
- Monthly billing: $[X]/month
- Annual billing: $[Y]/year (save $[Z] vs monthly)

### Expected ROI
[Quantify time savings, efficiency gains, or revenue impact where possible, using industry benchmarks from market research WITH SOURCES. Be realistic and conservative.]

**ROI Calculation:**
- Time saved per week: [X] hours
- Value per hour: $[Y]
- Monthly value: $[Z]
- Break-even point: [X] months

### Cost Optimization Strategies
[Suggest ways to reduce costs or maximize value]
- Annual payment discounts
- Team vs individual plans
- Free tier alternatives for testing
- Bundled vs a la carte pricing

---

## RISK ASSESSMENT & MITIGATION

### Potential Challenges
[List 3-5 realistic challenges the client might face, informed by industry best practices and market research]

### Mitigation Strategies
[For each challenge, provide specific actionable mitigation approach]

### Success Factors
[List critical factors for successful implementation]

---

## ALTERNATIVE SCENARIOS

### Budget-Constrained Alternative
[If budget is reduced by 30-50%, what would you recommend?]

**Modified Stack (Reduced Budget):**
- Total monthly cost: $[X]
- Key tradeoffs: [What you lose]
- Still achieves: [Core objectives maintained]

### Premium Alternative
[If budget allows 50-100% more investment, what enhanced capabilities could be added?]

**Enhanced Stack (Increased Budget):**
- Total monthly cost: $[X]
- Additional capabilities: [What you gain]
- Recommended for: [When premium makes sense]

### Different Scale Scenario
[How would recommendations change for 2x team size or 2x content volume?]

---

## SUPPORTING DATA

### Tool Comparison Matrix
[Create a simple text-based comparison of recommended tools vs alternatives on key criteria]

| Criteria | [Recommended Tool 1] | [Alternative] | Winner |
|----------|---------------------|---------------|--------|
| Price | $XX/month | $XX/month | [Tool] |
| Features | [Key features] | [Key features] | [Tool] |
| Ease of Use | [Rating] | [Rating] | [Tool] |
| Integration | [Rating] | [Rating] | [Tool] |

### Industry Benchmarks
[Include relevant statistics from market research about typical costs, implementation times, adoption rates - WITH SOURCES]

### Reference Documentation
[List official documentation, case studies, or resources for each recommended tool]

**Tool Resources:**
- **[Tool 1]:** [Official docs URL], [Getting started guide URL]
- **[Tool 2]:** [Official docs URL], [Getting started guide URL]

---

## MARKET RESEARCH & CONTEXT

**Industry Overview:**
Use web_search to find current market data related to the user's specific use case. For EVERY statistic or claim, include the source in parentheses.

Example format:
"The AI-powered video production market reached $2.8 billion in 2024 and is projected to grow at 22% CAGR through 2028 (MarketsandMarkets, 2024). According to Gartner's 2024 State of Marketing survey, 67% of marketing teams now use AI video tools, up from 34% in 2023."

Include:
- Market size and growth rate (with source)
- Current adoption trends and statistics (with source)
- Industry ROI benchmarks - be conservative (with source)
- Key challenges organizations typically face (with source)

**Competitive Landscape:**
Brief overview of how AI is currently transforming this specific area and what solutions are gaining traction (cite sources where applicable).

**Sources:**
After the Market Research section, include a "Sources" subsection listing all references:
- Source Name (Year): Brief description or URL if available
- Format as a simple bulleted list

Example:
**Sources:**
- MarketsandMarkets (2024): AI Video Production Market Report
- Gartner (2024): State of Marketing Survey
- Content Marketing Institute (2024): Industry Benchmark Report
- HubSpot Research (2024): Marketing Trends Report
- Forrester (2024): ROI Analysis

---

CRITICAL SOURCE CITATION REQUIREMENTS:
- EVERY statistic in Market Research section MUST include source in parentheses: (Source Name, Year)
- Include a "Sources:" subsection after Market Research with full references
- Use credible sources: Gartner, Forrester, IDC, McKinsey, industry associations, academic research
- If you cannot find a credible source for a claim, do not include the claim
- Industry benchmarks in ROI section should also cite sources when possible
- Prefer recent sources (2023-2024) over older data

TONE GUIDELINES:
- Professional and consultative, not casual or chatty
- Data-driven with specific numbers and CITED sources
- Balanced and objective, acknowledging tradeoffs
- Accessible to both technical and non-technical readers
- Conservative on ROI projections (under-promise, over-deliver)
- Analyst Note should be warm and welcoming while maintaining professionalism
- **NO EMOJIS** - Maintain professional consulting-firm tone throughout [NEW - #7]
- Use formal section headers without decorative elements
- Focus on substance over style

DATABASE:
{json.dumps(database['ai_tools'], indent=2)}"""
                
                status_text.text("📊 Researching market data... (30%)")
                progress_bar.progress(30)

                status_text.text("🤖 Generating strategic recommendations... (50%)")
                progress_bar.progress(50)
                
                # Retry logic for API overload errors
                max_retries = 3
                retry_count = 0
                message = None
                
                while retry_count < max_retries:
                    try:
                        if retry_count > 0:
                            status_text.text(f"⏳ API busy, retrying ({retry_count}/{max_retries})... (50%)")
                        
                        message = client.messages.create(
                            model="claude-sonnet-4-20250514",
                            max_tokens=4000,
                            tools=[{"type": "web_search_20250305", "name": "web_search"}],
                            system=[{
                                "type": "text",
                                "text": system_prompt,
                                "cache_control": {"type": "ephemeral"}
                            }],
                            messages=[{"role": "user", "content": user_query}]
                        )
                        break  # Success! Exit retry loop
                        
                    except anthropic.APIError as e:
                        if "overloaded" in str(e).lower() and retry_count < max_retries - 1:
                            retry_count += 1
                            wait_time = 5 * retry_count  # 5s, 10s, 15s
                            status_text.text(f"⏳ API overloaded, waiting {wait_time}s before retry {retry_count}/{max_retries}...")
                            import time
                            time.sleep(wait_time)
                        else:
                            # Not an overload error, or out of retries
                            raise
                
                if message is None:
                    raise Exception("Failed to generate report after multiple retries")
                
                status_text.text("✨ Finalizing report... (80%)")
                progress_bar.progress(80)
                
                # Extract text from response (handling tool use blocks)
                response_text = ""
                for block in message.content:
                    if hasattr(block, 'text'):
                        response_text += block.text
                
                # Clean response text - remove artifacts but PRESERVE mermaid blocks
                # Remove file paths
                response_text = re.sub(r'/[^\s]+\.py', '', response_text)
                # Remove thinking tags
                response_text = re.sub(r'<think>.*?</think>', '', response_text, flags=re.DOTALL)
                # Fix currency markdown
                response_text = response_text.replace('\\$', '$')
                
                # Store in session state so it persists after download
                st.session_state.generated_report = response_text
                st.session_state.report_query = user_query
                
                status_text.text("✅ Report complete! (100%)")
                progress_bar.progress(100)
                
                # Clear progress indicators after brief pause
                import time
                time.sleep(0.5)
                progress_bar.empty()
                status_text.empty()
                
                # Log query for analytics
                log_query(
                    user_query=user_query,
                    recommended_tools=None,
                    response_time=None
                )
                
                # Increment usage counter
                new_count = increment_usage()
                
                # Show success message
                remaining = 3 - new_count
                if remaining > 0:
                    st.success(f"✅ Report generated successfully! You have {remaining} reports remaining this month.")
                else:
                    st.info("✅ Report generated! This was your last free report this month. Resets on " + 
                           next_reset.strftime('%B %d, %Y'))
                
                # Report will display below via persistent display section
                
            except anthropic.APIError as e:
                if "overloaded" in str(e).lower():
                    st.error("⚠️ Anthropic API is experiencing high traffic")
                    st.info("💡 The API is temporarily overloaded. Please try again in 1-2 minutes. Your query has been saved.")
                    st.markdown(f"**Your query:** {user_query[:200]}...")
                elif "credit balance" in str(e).lower():
                    st.error(f"API Error: {str(e)}")
                    st.info("Please add credits to your Anthropic account at console.anthropic.com")
                else:
                    st.error(f"API Error: {str(e)}")
                    st.info("💡 Please try again in a moment.")
            except AttributeError as e:
                st.error("Response parsing error - this is usually due to tool use blocks in the API response")
                st.info("💡 The response format was unexpected. Please try again. The issue has been logged.")
            except Exception as e:
                st.error(f"Unexpected error: {str(e)}")
                st.info("💡 Please try again. If the issue persists, contact support.")

# ============================================================================
# DISPLAY REPORT (PERSISTS AFTER DOWNLOAD)
# ============================================================================

# Display report if it exists in session state (survives download button clicks)
if 'generated_report' in st.session_state and st.session_state.generated_report:
    response_text = st.session_state.generated_report
    user_query = st.session_state.get('report_query', '')
    
    # Display the report
    st.markdown('<div class="whitepaper-section">', unsafe_allow_html=True)
    st.markdown('<div class="whitepaper-header">Strategic Advisory Report</div>', unsafe_allow_html=True)
    
    # Quality validation
    quality_score, quality_issues, quality_warnings = validate_report_quality(response_text)
    st.markdown(display_quality_badge(quality_score), unsafe_allow_html=True)
    
    # Show quality feedback if there are issues
    if quality_issues or quality_warnings:
        with st.expander("📊 Quality Report Details", expanded=False):
            if quality_issues:
                st.markdown("**⚠️ Critical Issues:**")
                for issue in quality_issues:
                    st.markdown(f"- {issue}")
            
            if quality_warnings:
                st.markdown("**💡 Suggestions for Improvement:**")
                for warning in quality_warnings:
                    st.markdown(f"- {warning}")
            
            st.info("💡 The report has been generated, but addressing these items would improve quality.")
    
    # Display report content
    render_report_with_mermaid(response_text)
    st.markdown('</div>', unsafe_allow_html=True)
    
    # Export options
    st.markdown("---")
    col1, col2 = st.columns(2)
    
    with col1:
        # PowerPoint export
        try:
            pptx_file = generate_professional_pptx(response_text, user_query)
            st.download_button(
                label="📊 Export as Presentation (.pptx)",
                data=pptx_file,
                file_name=f"bulwise_advisory_{datetime.now().strftime('%Y%m%d')}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
                key="download_pptx"
            )
        except Exception as e:
            st.error(f"PowerPoint export error: {str(e)}")
    
    with col2:
        # PDF export
        try:
            pdf_file = generate_professional_pdf(response_text, user_query)
            if pdf_file and PDF_AVAILABLE:
                st.download_button(
                    label="📄 Export as PDF",
                    data=pdf_file,
                    file_name=f"bulwise_advisory_{datetime.now().strftime('%Y%m%d')}.pdf",
                    mime="application/pdf",
                    use_container_width=True,
                    key="download_pdf"
                )
            else:
                st.button("📄 Export as PDF (Installing...)", disabled=True, use_container_width=True, key="download_pdf_disabled")
        except Exception as e:
            st.button("📄 Export as PDF (Error)", disabled=True, use_container_width=True, key="download_pdf_error")
            st.caption(f"PDF generation error: {str(e)}")
    
    # Feedback section
    st.markdown("---")
    st.markdown("### 💬 Quick Feedback")
    
    col1, col2, col3 = st.columns([2, 1, 1])
    
    with col1:
        st.markdown("**Was this report helpful?**")
    
    with col2:
        if st.button("👍 Yes, helpful", key="feedback_yes_persist", use_container_width=True):
            log_feedback(user_query, "positive", None)
            st.success("Thank you for your feedback!")
            st.session_state.feedback_given = True
    
    with col3:
        if st.button("👎 Not helpful", key="feedback_no_persist", use_container_width=True):
            st.session_state.show_feedback_form = True
    
    # Feedback form
    if st.session_state.get('show_feedback_form', False) and not st.session_state.get('feedback_given', False):
        feedback_text = st.text_area(
            "What could we improve?",
            placeholder="Tell us what was missing or could be better...",
            key="feedback_details_persist"
        )
        
        if st.button("Submit Feedback", type="primary", key="submit_feedback_persist"):
            log_feedback(user_query, "negative", feedback_text)
            st.success("Thank you for your feedback! We'll use this to improve.")
            st.session_state.feedback_given = True
            st.session_state.show_feedback_form = False
    
    # Generate another report button
    if st.button("📝 Generate Another Report", type="primary", key="generate_another"):
        st.session_state.workflow_step = 1
        st.session_state.initial_query = ''
        st.session_state.clarifying_answers = {}
        if 'generated_report' in st.session_state:
            del st.session_state.generated_report
        if 'report_query' in st.session_state:
            del st.session_state.report_query
        if 'feedback_given' in st.session_state:
            del st.session_state.feedback_given
        if 'show_feedback_form' in st.session_state:
            del st.session_state.show_feedback_form
        st.rerun()
    
    # Consultancy CTA
    st.markdown("---")
    st.markdown("""
    <div style='background: linear-gradient(135deg, #2E7D32 0%, #1B5E20 100%); padding: 2rem; border-radius: 12px; text-align: center; margin: 2rem 0;'>
        <h3 style='color: white; margin-top: 0;'>Need Implementation Support?</h3>
        <p style='color: #E8F5E9; font-size: 1.1rem; margin-bottom: 1.5rem;'>
            Get personalized 1:1 consultation to help you implement these recommendations
        </p>
        <a href='mailto:hello@bulwise.io?subject=Implementation%20Consultation%20Request' 
           style='background: white; color: #2E7D32; padding: 0.75rem 2rem; border-radius: 8px; 
                  text-decoration: none; font-weight: 600; display: inline-block;'>
            📞 Book a Consultation
        </a>
    </div>
    """, unsafe_allow_html=True)

# Footer
st.markdown("---")
st.markdown("""
    <div style='text-align: center; color: #666; padding: 2rem 0;'>
        <p>Bulwise - AI Implementation Intelligence</p>
        <p>Questions? <a href='mailto:hello@bulwise.io'>hello@bulwise.io</a></p>
    </div>
""", unsafe_allow_html=True)
